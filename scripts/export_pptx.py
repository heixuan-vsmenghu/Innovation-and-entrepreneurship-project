# -*- coding: utf-8 -*-
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "ppt-outline" / "slides_full_content.md"
OUTPUT = ROOT / "deliverables" / "创策云UniPlanAI课堂汇报.pptx"
SCRIPT_OUTPUT = ROOT / "deliverables" / "创策云UniPlanAI汇报讲稿.md"

SLIDE_W = 13.333
SLIDE_H = 7.5

BLUE = RGBColor(23, 111, 193)
DEEP_BLUE = RGBColor(11, 47, 87)
CYAN = RGBColor(23, 162, 184)
DARK = RGBColor(18, 36, 60)
MID = RGBColor(75, 97, 120)
LIGHT_BG = RGBColor(239, 246, 252)
PALE_BLUE = RGBColor(231, 242, 251)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(18, 160, 113)
YELLOW = RGBColor(246, 180, 80)
RED = RGBColor(213, 85, 85)
LINE = RGBColor(196, 218, 238)


DEFAULT_SLIDES = [
    {
        "title": "项目标题",
        "bullets": ["创策云 UniPlan AI", "面向大学生创新创业项目的 AI 协作与商业计划书平台", "课程名称：大学生创新创业与就业指导", "团队成员：待补充"],
        "speaker": "大家好，我们小组汇报的项目是创策云 UniPlan AI。这个项目来源于我们完成创新创业课程作业时遇到的真实问题：小组项目从选题、调研到计划书和 PPT，流程很多，但工具很分散，所以我们希望设计一个更适合大学生团队的 AI 协作平台。",
    },
    {
        "title": "项目来源",
        "bullets": ["来源于大学生小组作业、创新创业课程和竞赛项目", "典型任务：选题、问卷、资料整理、商业计划书、PPT、课堂汇报", "当前问题：工具很多，但流程不连贯"],
        "speaker": "我们选择这个方向，是因为创新创业课程并不是只写一份文档，而是一个完整项目过程。很多小组会同时使用聊天群、在线文档、AI 工具和 PPT 软件，但这些工具之间缺少统一管理，最后经常变成临时拼接。",
    },
    {
        "title": "用户痛点",
        "bullets": ["选题难：不知道专业方向如何转化为创业项目", "资料散：政策、竞品、问卷、访谈分散在不同地方", "分工乱：成员任务不清楚，进度靠群消息催", "写作难：商业计划书结构复杂，市场和财务部分难写", "版本乱：多人修改后容易不一致"],
        "speaker": "我们总结了五类痛点。最突出的是选题、资料、分工、写作和版本。尤其是商业计划书，不少同学不是没有想法，而是不知道应该按照什么结构展开，哪些地方需要数据，哪些地方需要分析。",
    },
    {
        "title": "调研设计",
        "bullets": ["问卷主题：大学生小组项目协作与商业计划书写作痛点调研", "对象：在校大学生，尤其是参加过小组项目的学生", "模块：基本信息、协作痛点、写作困难、AI 接受度、付费意愿、开放建议", "样本数量：【待问卷调研后补充】"],
        "speaker": "为了验证需求，我们设计了一份 15 题以内的问卷，后续会通过班级群和同学社群发放。问卷不会直接假设用户一定需要我们的产品，而是先了解他们在项目协作和计划书写作中的真实困难。",
    },
    {
        "title": "目标用户",
        "bullets": ["创新创业课程小组", "创新创业竞赛团队", "大学生创新创业训练计划团队", "指导教师、助教和创新创业中心"],
        "speaker": "创策云早期不面向所有办公用户，而是聚焦高校创新创业项目。学生团队是直接用户，教师和助教是重要影响者。如果平台能帮助学生规范过程，也能减轻教师查看材料和指导项目的压力。",
    },
    {
        "title": "产品定位",
        "bullets": ["一句话定位：面向大学生创新创业项目的 AI 协作平台", "核心价值：把选题、调研、写作、PPT 和分工放在一个工作台", "使用原则：AI 辅助，不替代真实调研和人工判断"],
        "speaker": "我们的定位不是简单的 AI 代写工具，而是项目协作平台。AI 可以帮助生成大纲和草稿，但真实数据、团队判断和最终修改仍然由学生完成。这样既提高效率，也符合课程学习目标。",
    },
    {
        "title": "核心功能",
        "bullets": ["选题分析：从专业方向生成项目建议", "资料整理：集中管理政策、竞品和访谈", "问卷调研：生成问卷与分析模板", "商业计划书生成：按章节输出草稿和提示", "PPT 大纲：生成页面内容与讲稿", "任务看板：记录分工、进度和负责人"],
        "speaker": "平台的功能围绕项目全过程设计。比如商业计划书生成页不仅生成文字，还会提醒哪些地方需要问卷数据；任务看板不仅显示任务，还能让课堂展示时看出小组成员各自负责什么。",
    },
    {
        "title": "原型展示",
        "bullets": ["已完成 React + Vite 前端原型", "包含：首页、功能介绍、项目工作台、商业计划书生成页、任务看板页", "静态模拟数据，不接真实 AI API"],
        "speaker": "这一页可以切换到我们的网页原型。首页展示项目进度和交付物，工作台展示资料库和版本记录，计划书生成页展示章节草稿，任务看板展示小组分工。当前原型重点是课堂展示可信度，暂时不接真实 AI。",
    },
    {
        "title": "商业计划书支持",
        "bullets": ["按教材创业计划书格式组织章节", "支持执行概要、公司概述、产品服务、市场分析、营销计划、财务风险等内容", "对真实数据位置标注：【待问卷调研后补充】"],
        "speaker": "商业计划书是课程提交的重点，所以我们按照教材格式整理了正文结构。平台会把大标题和小标题先搭好，再帮助学生填充内容。对于市场规模、转化率、收入预测等不能凭空编造的数据，会统一标注待调研后补充。",
    },
    {
        "title": "商业模式",
        "bullets": ["免费基础版：项目创建、基础大纲、任务看板", "学生会员版：更多模板、更多 AI 生成次数、版本记录", "高校课程版：班级项目管理、材料导出、教师反馈", "竞赛服务包：路演问答、项目诊断、PPT 优化"],
        "speaker": "商业模式上，我们采用从免费基础版到增值服务的思路。早期先让学生愿意使用，后续再考虑高级模板、竞赛服务和高校课程版。具体价格和转化率需要等问卷和试用结果出来后再判断。",
    },
    {
        "title": "营销计划",
        "bullets": ["校园冷启动：从本课程和同专业小组试用开始", "内容推广：分享计划书模板、问卷模板、PPT 结构", "教师合作：争取课程教师或助教推荐", "竞赛节点：围绕报名和路演阶段推广"],
        "speaker": "营销上我们不会一开始做大范围投放，而是从校园场景开始。因为这个项目的用户就在我们身边，课程小组、竞赛团队和教师都能给到直接反馈。先验证真实需求，再考虑扩大推广。",
    },
    {
        "title": "团队分工",
        "bullets": ["项目负责人：统筹方向与进度", "市场调研负责人：问卷、访谈和数据分析", "产品与原型负责人：功能流程和前端原型", "商业计划书负责人：正文撰写和格式整理", "PPT 与汇报负责人：视觉设计和课堂讲稿", "财务与风险分析负责人：财务测算和风险应对"],
        "speaker": "我们按照课程作业和创业团队组建的要求进行了分工。每个人都有明确模块，后续可以把学号和姓名补充到团队分工表中。这样也方便打印材料和课堂展示。",
    },
    {
        "title": "财务与融资",
        "bullets": ["课程阶段：不进行真实融资，以原型和调研验证为主", "成本：开发时间、服务器、AI 调用、推广物料", "收入：学生会员、高校课程版、竞赛服务包", "财务数据：【待问卷调研后补充】"],
        "speaker": "财务方面我们采用保守估计。课程阶段重点不是赚钱，而是验证需求。后续如果用户反馈较好，可以申请校级创新创业训练项目或创业基金，再进行更完整的成本和收入测算。",
    },
    {
        "title": "风险应对",
        "bullets": ["内容真实性风险：AI 结果需人工核验", "教育合规风险：强调辅助学习，不替代作业", "付费意愿风险：提供免费基础版", "竞争风险：聚焦高校创新创业场景", "团队执行风险：用分工表和时间计划管理"],
        "speaker": "我们也考虑了风险。最重要的是 AI 内容不能直接当成真实结论，所以平台要提醒用户补充问卷和访谈数据。同时，产品不能被定位成代写工具，而应是帮助学生理解结构、推进协作的学习辅助工具。",
    },
    {
        "title": "发展规划与总结",
        "bullets": ["第一阶段：完成项目骨架和原型", "第二阶段：完善作业包、问卷和 PPT", "第三阶段：校内小范围试用", "第四阶段：接入真实 AI 和教师端", "总结：让大学生创新创业项目更清晰、更协作、更可信"],
        "speaker": "最后总结一下，创策云 UniPlan AI 解决的是大学生创新创业项目中的真实协作和写作问题。我们希望它不是一个看起来很大的概念，而是一个从课程作业出发、能实际帮助小组推进项目的平台。谢谢大家。",
    },
]


def parse_slides(markdown):
    pattern = re.compile(r"^## 第\s*(\d+)\s*页[：:](.+)$", re.MULTILINE)
    matches = list(pattern.finditer(markdown))
    if len(matches) != 15:
        return DEFAULT_SLIDES

    slides = []
    for idx, match in enumerate(matches):
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(markdown)
        block = markdown[start:end].strip()
        title = match.group(2).strip()
        bullets = []
        speaker = ""
        mode = None
        for raw in block.splitlines():
            line = raw.strip()
            if not line:
                continue
            if line.startswith("页面文字"):
                mode = "bullets"
                continue
            if line.startswith("图示建议"):
                mode = "visual"
                continue
            if line.startswith("汇报讲稿"):
                mode = "speaker"
                speaker = line.split("：", 1)[-1].strip() if "：" in line else ""
                continue
            if mode == "bullets" and line.startswith("- "):
                bullets.append(line[2:].strip())
            elif mode == "speaker":
                speaker = (speaker + " " + line).strip()
        fallback = DEFAULT_SLIDES[idx]
        slides.append(
            {
                "title": title or fallback["title"],
                "bullets": (bullets or fallback["bullets"])[:6],
                "speaker": speaker or fallback["speaker"],
            }
        )
    return slides


def set_text(paragraph, text, size=18, bold=False, color=DARK, align=None):
    paragraph.text = ""
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=None, valign=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    if valign:
        frame.vertical_anchor = valign
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    set_text(frame.paragraphs[0], text, size=size, bold=bold, color=color, align=align)
    return shape


def add_shape(slide, shape_type, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_card(slide, x, y, w, h, title, body="", fill=WHITE, accent=BLUE, title_size=18, body_size=13):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=LINE)
    stripe = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, fill=accent, line=accent)
    add_text(slide, x + 0.18, y + 0.16, w - 0.32, 0.34, title, title_size, True, DARK)
    if body:
        add_text(slide, x + 0.18, y + 0.62, w - 0.32, h - 0.78, body, body_size, False, MID)
    return card


def add_title(slide, title, number=None):
    add_text(slide, 0.65, 0.45, 10.8, 0.55, title, 28, True, DARK)
    if number is not None:
        add_text(slide, 11.9, 0.52, 0.8, 0.3, f"{number:02d}", 13, True, BLUE, PP_ALIGN.RIGHT)


def add_footer(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.16, SLIDE_W, 0.09, fill=PALE_BLUE, line=PALE_BLUE)
    add_text(slide, 0.68, 7.02, 4.2, 0.25, "创策云 UniPlan AI", 10, True, BLUE)


def apply_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.18, fill=BLUE, line=BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.18, SLIDE_W, 0.03, fill=CYAN, line=CYAN)
    add_footer(slide)


def add_bullets(slide, bullets, x=0.8, y=1.35, w=5.9, h=4.8, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets[:5]):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = DARK


def slide_cover(slide, item):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP_BLUE
    for x, y, w, h, color in [
        (8.7, 0.55, 4.4, 0.9, BLUE),
        (9.55, 1.72, 3.1, 0.18, CYAN),
        (0.0, 6.72, 13.33, 0.28, BLUE),
        (0.0, 7.0, 13.33, 0.12, CYAN),
    ]:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=color, line=color)
    for i in range(7):
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.75 + i * 0.48, 0.75 + i * 0.28, 0.22, 0.22, fill=RGBColor(64, 135, 203), line=RGBColor(64, 135, 203))
    add_text(slide, 0.9, 1.55, 10.8, 0.75, "创策云 UniPlan AI", 40, True, WHITE)
    add_text(slide, 0.95, 2.48, 10.1, 0.55, "面向大学生创新创业项目的 AI 协作与商业计划书平台", 23, False, RGBColor(220, 238, 255))
    add_text(slide, 0.98, 3.55, 7.5, 0.38, "课程名称：大学生创新创业与就业指导", 17, False, RGBColor(210, 229, 247))
    add_text(slide, 0.98, 4.02, 7.5, 0.38, "团队成员：待补充", 17, False, RGBColor(210, 229, 247))
    add_text(slide, 0.98, 5.45, 8.4, 0.42, "让大学生创新创业项目更清晰、更协作、更可信", 19, True, WHITE)
    hub = add_shape(slide, MSO_SHAPE.OVAL, 9.15, 3.08, 2.15, 2.15, fill=WHITE, line=CYAN)
    add_text(slide, 9.48, 3.77, 1.52, 0.42, "AI\n协作", 23, True, BLUE, PP_ALIGN.CENTER)


def slide_process(slide, item):
    add_title(slide, "项目来源：从课程任务到完整交付", 2)
    steps = ["选题", "调研", "计划书", "PPT", "汇报"]
    for i, step in enumerate(steps):
        x = 0.75 + i * 2.42
        add_shape(slide, MSO_SHAPE.RIGHT_ARROW if i < 4 else MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, 1.9, 1.0, fill=BLUE if i < 4 else CYAN, line=BLUE if i < 4 else CYAN)
        add_text(slide, x + 0.18, 2.37, 1.3, 0.36, step, 22, True, WHITE, PP_ALIGN.CENTER)
    add_bullets(slide, item["bullets"], x=1.0, y=4.05, w=10.9, h=1.7, size=19)


def slide_pain(slide, item):
    add_title(slide, "用户痛点：小组项目常见阻力", 3)
    cards = [
        ("选题难", "专业方向难转成项目"),
        ("资料散", "政策、竞品、问卷分散"),
        ("分工乱", "任务边界和进度不清"),
        ("写作难", "计划书结构复杂"),
        ("版本乱", "多人修改口径不一致"),
    ]
    for i, (title, body) in enumerate(cards):
        add_card(slide, 0.65 + i * 2.52, 2.05, 2.12, 2.35, title, body, fill=LIGHT_BG, accent=[BLUE, CYAN, YELLOW, RED, GREEN][i], title_size=21, body_size=14)
    add_text(slide, 1.0, 5.35, 11.1, 0.55, "痛点集中在“从空白开始”和“多人协作失控”两个环节。", 22, True, DARK, PP_ALIGN.CENTER)


def slide_survey(slide, item):
    add_title(slide, "调研设计：先验证真实需求", 4)
    modules = ["基本信息", "协作痛点", "写作困难", "AI 接受度", "付费意愿"]
    for i, module in enumerate(modules):
        x = 0.85 + i * 2.35
        add_card(slide, x, 2.0, 1.85, 1.28, module, "", fill=PALE_BLUE, accent=BLUE, title_size=17)
        if i < len(modules) - 1:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x + 1.74, 2.43, 0.7, 0.36, fill=CYAN, line=CYAN)
    add_bullets(slide, item["bullets"], x=1.05, y=4.15, w=10.6, h=1.55, size=18)


def slide_users(slide, item):
    add_title(slide, "目标用户：聚焦高校项目团队", 5)
    users = [
        ("课程小组", "完成课堂作业与计划书"),
        ("竞赛团队", "准备路演和项目材料"),
        ("训练计划团队", "沉淀长期项目过程"),
        ("教师 / 助教", "查看过程与指导项目"),
    ]
    for i, (title, body) in enumerate(users):
        add_card(slide, 1.0 + i * 3.0, 2.0, 2.45, 2.55, title, body, fill=WHITE, accent=[BLUE, CYAN, GREEN, YELLOW][i], title_size=19, body_size=15)
    add_text(slide, 1.15, 5.35, 10.8, 0.42, "早期主攻学生团队，后续延伸到教师端和课程版。", 21, True, BLUE, PP_ALIGN.CENTER)


def slide_position(slide, item):
    add_title(slide, "产品定位：一站式 AI 协作工作台", 6)
    add_shape(slide, MSO_SHAPE.OVAL, 5.42, 2.62, 2.3, 2.3, fill=BLUE, line=BLUE)
    add_text(slide, 5.72, 3.28, 1.7, 0.55, "创策云", 24, True, WHITE, PP_ALIGN.CENTER)
    nodes = [
        ("选题", 2.1, 1.55), ("调研", 5.4, 1.1), ("写作", 8.75, 1.55),
        ("PPT", 2.05, 5.25), ("看板", 5.38, 5.72), ("版本", 8.72, 5.25),
    ]
    for title, x, y in nodes:
        add_shape(slide, MSO_SHAPE.OVAL, x, y, 1.42, 0.72, fill=PALE_BLUE, line=CYAN)
        add_text(slide, x + 0.12, y + 0.16, 1.18, 0.3, title, 17, True, BLUE, PP_ALIGN.CENTER)
    add_text(slide, 0.95, 6.35, 11.4, 0.42, "AI 辅助结构化推进，真实调研与最终判断仍由学生完成。", 19, True, DARK, PP_ALIGN.CENTER)


def slide_features(slide, item):
    add_title(slide, "核心功能：覆盖项目全过程", 7)
    features = [
        ("选题分析", "从专业方向生成建议"),
        ("资料整理", "集中管理资料来源"),
        ("问卷调研", "生成问卷与分析模板"),
        ("计划书生成", "按章节输出草稿"),
        ("PPT 大纲", "生成页面和讲稿"),
        ("任务看板", "记录分工和进度"),
    ]
    for i, (title, body) in enumerate(features):
        row, col = divmod(i, 3)
        add_card(slide, 1.0 + col * 3.85, 1.65 + row * 2.18, 3.25, 1.55, title, body, fill=LIGHT_BG, accent=BLUE if i % 2 == 0 else CYAN, title_size=18, body_size=13)


def slide_prototype(slide, item):
    add_title(slide, "原型展示：五个页面可课堂演示", 8)
    pages = [
        ("首页", "项目进度 / 交付物"),
        ("功能页", "六类核心能力"),
        ("工作台", "资料库 / 版本记录"),
        ("计划书生成", "章节草稿 / 数据提醒"),
        ("任务看板", "待办 / 进行中 / 已完成"),
    ]
    for i, (title, body) in enumerate(pages):
        x = 0.65 + i * 2.52
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.65, 2.16, 3.6, fill=WHITE, line=LINE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.65, 2.16, 0.42, fill=BLUE, line=BLUE)
        add_text(slide, x + 0.16, 1.76, 1.84, 0.18, title, 11, True, WHITE, PP_ALIGN.CENTER)
        for j in range(3):
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.22, 2.35 + j * 0.55, 1.72, 0.16, fill=RGBColor(221, 234, 246), line=RGBColor(221, 234, 246))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.22, 4.15, 1.72, 0.48, fill=PALE_BLUE, line=PALE_BLUE)
        add_text(slide, x + 0.22, 5.02, 1.72, 0.4, body, 12, False, MID, PP_ALIGN.CENTER)
    add_text(slide, 0.9, 6.0, 11.6, 0.38, "静态模拟数据，重点展示项目流程和课堂可信度。", 18, True, DARK, PP_ALIGN.CENTER)


def slide_plan_tree(slide, item):
    add_title(slide, "商业计划书支持：按课程格式组织", 9)
    root_x, root_y = 0.95, 1.45
    add_card(slide, root_x, root_y, 2.4, 0.82, "创业计划书", "", fill=BLUE, accent=BLUE, title_size=18)
    chapters = ["执行概要", "公司概述", "产业分析", "市场与营销", "产品服务", "团队组织", "运营开发", "财务风险", "附录材料"]
    for i, chapter in enumerate(chapters):
        row, col = divmod(i, 3)
        x = 4.0 + col * 2.65
        y = 1.25 + row * 1.35
        add_card(slide, x, y, 2.05, 0.78, chapter, "", fill=LIGHT_BG, accent=CYAN, title_size=15)
    add_text(slide, 0.95, 5.82, 11.3, 0.55, "所有真实数据位置保留【待问卷调研后补充】，避免凭空编造。", 20, True, BLUE, PP_ALIGN.CENTER)


def slide_business_model(slide, item):
    add_title(slide, "商业模式：九要素设计", 10)
    cells = [
        "客户细分", "价值主张", "渠道通路",
        "客户关系", "收入来源", "关键业务",
        "核心资源", "重要伙伴", "成本结构",
    ]
    for i, cell in enumerate(cells):
        row, col = divmod(i, 3)
        add_card(slide, 1.15 + col * 3.55, 1.35 + row * 1.62, 3.0, 1.05, cell, "", fill=PALE_BLUE if i % 2 == 0 else WHITE, accent=BLUE, title_size=17)
    add_text(slide, 1.3, 6.35, 10.8, 0.35, "先免费验证需求，再探索学生会员、课程版和竞赛服务包。", 18, True, DARK, PP_ALIGN.CENTER)


def slide_marketing(slide, item):
    add_title(slide, "营销计划：校园冷启动路径", 11)
    steps = [("课程试用", "本课程和同专业小组"), ("内容推广", "模板 / 问卷 / PPT 结构"), ("教师合作", "助教与课程推荐"), ("竞赛节点", "报名和路演阶段")]
    for i, (title, body) in enumerate(steps):
        x = 0.9 + i * 3.05
        add_shape(slide, MSO_SHAPE.CHEVRON, x, 2.25, 2.55, 1.15, fill=[BLUE, CYAN, GREEN, YELLOW][i], line=[BLUE, CYAN, GREEN, YELLOW][i])
        add_text(slide, x + 0.22, 2.45, 1.85, 0.32, title, 17, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, 3.72, 2.1, 0.58, body, 14, False, MID, PP_ALIGN.CENTER)
    add_bullets(slide, item["bullets"], x=1.15, y=5.0, w=10.8, h=1.3, size=17)


def slide_team(slide, item):
    add_title(slide, "团队分工：目标一致、能力互补", 12)
    add_card(slide, 5.2, 1.25, 2.8, 0.85, "项目负责人", "统筹方向与进度", fill=BLUE, accent=BLUE, title_size=18, body_size=12)
    roles = [
        ("市场调研", "问卷 / 访谈 / 数据", 0.85, 2.95),
        ("产品与原型", "功能流程 / 前端原型", 3.95, 2.95),
        ("计划书", "正文 / 格式 / 打印", 7.05, 2.95),
        ("PPT 与汇报", "视觉 / 讲稿 / 演练", 10.15, 2.95),
        ("财务与风险", "测算 / 融资 / 风险", 5.5, 5.0),
    ]
    for title, body, x, y in roles:
        add_card(slide, x, y, 2.35, 1.05, title, body, fill=LIGHT_BG, accent=CYAN, title_size=16, body_size=12)
    add_text(slide, 0.9, 6.42, 11.7, 0.3, "后续公司治理按贡献、投入和责任合理安排，不虚构股权比例。", 16, True, MID, PP_ALIGN.CENTER)


def slide_finance(slide, item):
    add_title(slide, "财务与融资：先验证，再投入", 13)
    add_card(slide, 1.0, 1.55, 5.2, 4.2, "成本结构", "开发时间\n服务器与域名\nAI 调用\n推广物料\n调研活动", fill=LIGHT_BG, accent=RED, title_size=21, body_size=20)
    add_card(slide, 7.1, 1.55, 5.2, 4.2, "收入来源", "学生会员\n高校课程版\n竞赛服务包\n模板增值服务\n教师端服务", fill=LIGHT_BG, accent=GREEN, title_size=21, body_size=20)
    add_text(slide, 1.1, 6.22, 11.0, 0.4, "课程阶段不做真实融资，财务数据待问卷和试用后补充。", 18, True, BLUE, PP_ALIGN.CENTER)


def slide_risk(slide, item):
    add_title(slide, "风险应对：识别、处理、监控", 14)
    rows = [
        ("内容真实性", "人工核验 + 数据占位"),
        ("教育合规", "辅助学习 + 过程留痕"),
        ("付费意愿", "免费基础版验证"),
        ("竞争风险", "聚焦高校课程场景"),
        ("团队执行", "分工表 + 时间计划"),
    ]
    x0, y0 = 1.0, 1.45
    add_shape(slide, MSO_SHAPE.RECTANGLE, x0, y0, 11.3, 0.5, fill=BLUE, line=BLUE)
    add_text(slide, x0 + 0.25, y0 + 0.11, 4.0, 0.22, "主要风险", 13, True, WHITE)
    add_text(slide, x0 + 5.55, y0 + 0.11, 4.0, 0.22, "解决方案", 13, True, WHITE)
    for i, (risk, action) in enumerate(rows):
        y = y0 + 0.62 + i * 0.75
        fill = LIGHT_BG if i % 2 == 0 else WHITE
        add_shape(slide, MSO_SHAPE.RECTANGLE, x0, y, 5.45, 0.62, fill=fill, line=LINE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x0 + 5.45, y, 5.85, 0.62, fill=fill, line=LINE)
        add_text(slide, x0 + 0.25, y + 0.16, 4.7, 0.22, risk, 14, True, DARK)
        add_text(slide, x0 + 5.7, y + 0.16, 5.1, 0.22, action, 14, False, DARK)


def slide_roadmap(slide, item):
    add_title(slide, "发展规划：从课程原型到校园平台", 15)
    stages = [
        ("第一阶段", "项目骨架和原型"),
        ("第二阶段", "作业包、问卷和 PPT"),
        ("第三阶段", "校内小范围试用"),
        ("第四阶段", "真实 AI 和教师端"),
    ]
    y = 3.15
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.35, y + 0.32, 10.65, 0.08, fill=LINE, line=LINE)
    for i, (stage, body) in enumerate(stages):
        x = 1.15 + i * 3.15
        add_shape(slide, MSO_SHAPE.OVAL, x, y, 0.72, 0.72, fill=[BLUE, CYAN, GREEN, YELLOW][i], line=[BLUE, CYAN, GREEN, YELLOW][i])
        add_text(slide, x + 0.07, y + 0.19, 0.58, 0.22, str(i + 1), 14, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, x - 0.55, y - 0.72, 1.85, 0.3, stage, 15, True, BLUE, PP_ALIGN.CENTER)
        add_text(slide, x - 0.75, y + 0.92, 2.15, 0.42, body, 14, False, DARK, PP_ALIGN.CENTER)
    add_text(slide, 1.0, 5.55, 11.3, 0.55, "让大学生创新创业项目更清晰、更协作、更可信", 24, True, DEEP_BLUE, PP_ALIGN.CENTER)


SLIDE_BUILDERS = [
    slide_cover,
    slide_process,
    slide_pain,
    slide_survey,
    slide_users,
    slide_position,
    slide_features,
    slide_prototype,
    slide_plan_tree,
    slide_business_model,
    slide_marketing,
    slide_team,
    slide_finance,
    slide_risk,
    slide_roadmap,
]


def build_ppt(slides_data):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    for index, item in enumerate(slides_data[:15]):
        slide = prs.slides.add_slide(blank)
        if index == 0:
            slide_cover(slide, item)
        else:
            apply_background(slide)
            SLIDE_BUILDERS[index](slide, item)
    return prs


def export_speaker_notes(slides_data):
    lines = ["# 创策云 UniPlan AI 汇报讲稿", "", "建议总时长：6—8 分钟。", ""]
    for idx, item in enumerate(slides_data[:15], start=1):
        lines.append(f"## 第 {idx} 页：{item['title']}")
        lines.append("")
        lines.append(item["speaker"])
        lines.append("")
    SCRIPT_OUTPUT.write_text("\n".join(lines), encoding="utf-8")


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    markdown = SOURCE.read_text(encoding="utf-8")
    slides_data = parse_slides(markdown)
    prs = build_ppt(slides_data)
    prs.save(OUTPUT)
    export_speaker_notes(slides_data)
    print(f"exported {OUTPUT}")
    print(f"exported {SCRIPT_OUTPUT}")


if __name__ == "__main__":
    main()
